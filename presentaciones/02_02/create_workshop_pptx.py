"""
Script para generar la presentación del Workshop SEMIL-USP
Basado en SLIDES_MEJORADOS_WORKSHOP.md
Autor: Science Team
Fecha: 02/02/2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Colores del proyecto (usando RGBColor correctamente)
COLORS = {
    'governance': RGBColor(0x2E, 0x86, 0xAB),  # Azul
    'biodiversity': RGBColor(0x28, 0xA7, 0x45),  # Verde
    'climate': RGBColor(0xF5, 0x7C, 0x00),  # Naranja
    'health': RGBColor(0xDC, 0x35, 0x45),  # Rojo
    'vulnerability': RGBColor(0x6F, 0x42, 0xC1),  # Morado
    'dark': RGBColor(0x2C, 0x3E, 0x50),  # Gris oscuro
    'light': RGBColor(0xEC, 0xF0, 0xF1),  # Gris claro
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'black': RGBColor(0x00, 0x00, 0x00),
}

def add_title_slide(prs, title, subtitle=""):
    """Agrega slide de título"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Fondo con color
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark']
    bg.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.5), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['light']
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(20)

    return slide

def add_content_slide(prs, title, content_items, highlight_color=None):
    """Agrega slide de contenido con bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Barra superior de color
    if highlight_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = highlight_color
        bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # Contenido
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.font.size = Pt(item.get('size', 18))
            p.font.bold = item.get('bold', False)
            if item.get('color'):
                p.font.color.rgb = item['color']
            else:
                p.font.color.rgb = COLORS['dark']
            p.level = item.get('level', 0)
        else:
            if item:
                p.text = f"  {item}"
            else:
                p.text = ""
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['dark']

        p.space_before = Pt(8)

    return slide

def add_table_slide(prs, title, headers, rows, highlight_color=None):
    """Agrega slide con tabla"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Barra superior
    if highlight_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = highlight_color
        bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # Tabla
    num_rows = len(rows) + 1
    num_cols = len(headers)

    # Calcular ancho de columnas
    table_width = Inches(12.5)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.3), table_width, Inches(0.4 * num_rows)).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['dark']
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(12)
        para.font.color.rgb = COLORS['white']
        para.alignment = PP_ALIGN.CENTER

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = COLORS['dark']
            para.alignment = PP_ALIGN.CENTER
            # Alternar colores de fila
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light']

    return slide

def add_highlight_slide(prs, main_text, subtext="", number="", color=None):
    """Agrega slide de highlight con texto grande"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color if color else COLORS['governance']
    bg.line.fill.background()

    # Número grande si existe
    if number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(1.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(96)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    # Texto principal
    y_pos = Inches(3) if number else Inches(2.5)
    main_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12.5), Inches(2))
    tf = main_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = main_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Subtexto
    if subtext:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.5), Inches(1.5))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_question_slide(prs, question, options=None):
    """Agrega slide de pregunta"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark']
    bg.line.fill.background()

    # Pregunta
    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.5), Inches(2))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Opciones
    if options:
        opt_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.5), Inches(2.5))
        tf = opt_box.text_frame
        tf.word_wrap = True
        for i, opt in enumerate(options):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = opt
            p.font.size = Pt(22)
            p.font.color.rgb = COLORS['light']
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(12)

    return slide

def create_presentation():
    """Crea la presentación completa"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =====================
    # BLOQUE A: EL PROBLEMA
    # =====================

    # SLIDE 1: Título
    add_title_slide(
        prs,
        "Riesgos climaticos, salud y perdida de biodiversidad\npredicen politicas de adaptacion...\nexcepto donde mas se necesitan",
        "Evidencia de 645 municipios de Sao Paulo (2010-2019)\n\nWorkshop SEMIL-USP | 24 de febrero 2026"
    )

    # SLIDE 2: Pregunta central
    add_question_slide(
        prs,
        "La gobernanza climatica en Sao Paulo\nes PROACTIVA o REACTIVA?",
        [
            "PROACTIVA: Anticipa riesgos, invierte antes de crisis",
            "REACTIVA: Responde a desastres, invierte despues de crisis",
            "",
            "Esta es la pregunta que guia nuestro analisis..."
        ]
    )

    # SLIDE 3: Marco Nexus
    add_content_slide(
        prs,
        "El Nexo Biodiversidad-Clima-Salud-Gobernanza",
        [
            {'text': "Cuatro dimensiones interconectadas:", 'size': 20, 'bold': True, 'color': COLORS['dark']},
            "",
            {'text': "BIODIVERSIDAD - Servicios ecosistemicos (polinizacion, regulacion)", 'size': 18},
            {'text': "   Efecto dilucion: ecosistemas diversos reducen transmision de patogenos", 'size': 14, 'color': COLORS['biodiversity']},
            "",
            {'text': "CLIMA - Riesgos (inundaciones, fuego, sequia, calor)", 'size': 18},
            {'text': "   Exposicion diferencial segun territorio y vulnerabilidad", 'size': 14, 'color': COLORS['climate']},
            "",
            {'text': "SALUD - Indicadores de impacto (dengue, malaria, cardiovascular)", 'size': 18},
            {'text': "   Proxy de la calidad del ambiente y acceso a servicios", 'size': 14, 'color': COLORS['health']},
            "",
            {'text': "GOBERNANZA - Capacidad adaptativa (UAI - 5 dimensiones)", 'size': 18},
            {'text': "   Previene riesgos o responde a ellos?", 'size': 14, 'color': COLORS['governance']},
        ],
        COLORS['governance']
    )

    # SLIDE 4: Nuestros datos
    add_table_slide(
        prs,
        "645 municipios x 104 variables x 10 anos",
        ["Dimension", "Variables", "Fuente"],
        [
            ["Gobernanza", "26 indicadores UAI", "Neder et al. 2021"],
            ["Biodiversidad", "Riqueza, cobertura, polinizacion", "IBGE, MapBiomas"],
            ["Clima", "Inundacion, fuego, estres hidrico", "INPE, ANA"],
            ["Salud", "7 enfermedades + mortalidad", "DATASUS"],
            ["Vulnerabilidad", "Pobreza, ruralidad, demografia", "CadUnico, IBGE"],
        ],
        COLORS['dark']
    )

    # SLIDE 5: Pregunta del workshop
    add_question_slide(
        prs,
        "Si tuvieramos recursos LIMITADOS\npara intervenir...\n\nComo priorizariamos los 645 municipios?",
        [
            "Donde hay mas riesgo climatico?",
            "Donde hay mas vulnerabilidad social?",
            "Donde hay menos gobernanza?",
            "",
            "Los datos revelan patrones que desafian nuestra intuicion"
        ]
    )

    # ==========================
    # BLOQUE B: LO QUE ENCONTRAMOS
    # ==========================

    # SLIDE 6: Hallazgo 1 - Paradoja gobernanza
    add_content_slide(
        prs,
        "Hallazgo 1: Mas gobernanza NO significa menos riesgo climatico",
        [
            {'text': "La correlacion es POSITIVA (r = +0.29)", 'size': 24, 'bold': True, 'color': COLORS['climate']},
            "",
            "Los municipios con MEJOR gobernanza tienen MAYOR riesgo climatico",
            "",
            {'text': "Por que? La gobernanza responde al riesgo, no lo previene.", 'size': 20, 'bold': True},
            "",
            {'text': "Modelo causal encontrado:", 'size': 18},
            "",
            {'text': "RIESGO --> DESASTRE --> INVERSION --> GOBERNANZA", 'size': 22, 'bold': True, 'color': COLORS['governance']},
            {'text': "(reactivo, no preventivo)", 'size': 16, 'color': COLORS['health']},
            "",
            "Los municipios desarrollan capacidad DESPUES de sufrir eventos adversos",
        ],
        COLORS['climate']
    )

    # SLIDE 7: Hallazgo 2 - Injusticia climática
    add_highlight_slide(
        prs,
        "La pobreza anula el efecto de la gobernanza",
        "27% de la varianza en gobernanza es explicada por pobreza\n\nLos municipios mas pobres tienen menos capacidad adaptativa\nPero son los que mas la necesitan.",
        "27%",
        COLORS['vulnerability']
    )

    # SLIDE 8: Hallazgo 3 - Efecto dilución
    add_table_slide(
        prs,
        "Hallazgo 3: La biodiversidad protege - pero no a todos igual",
        ["Relacion", "Correlacion", "Interpretacion"],
        [
            ["Bosque - Dengue", "r = -0.45", "Protege fuerte"],
            ["Bosque - Malaria", "r = -0.45", "Protege fuerte"],
            ["Bosque - Diarrea", "r = -0.41", "Protege moderado"],
            ["", "", ""],
            ["PERO...", "", ""],
            ["Alta pobreza", "-->", "Reduce beneficio"],
            ["Alta % pob. negra", "-->", "Reduce beneficio"],
        ],
        COLORS['biodiversity']
    )

    # SLIDE 9: Hallazgo 4 - Predictores por enfermedad
    add_table_slide(
        prs,
        "Hallazgo 4: Cada enfermedad tiene su propia firma de predictores",
        ["Enfermedad", "Mejor predictor", "Gobernanza ayuda?"],
        [
            ["Dengue", "% Rural (-)", "Si (indice general)"],
            ["Malaria", "Cobertura forestal (+)", "Si (UAI riesgo)"],
            ["Leptospirosis", "Biodiversidad (-)", "No significativo"],
            ["Leishmaniasis", "Estres hidrico (-)", "No significativo"],
            ["Diarrea", "Riesgo fuego (+)", "No significativo"],
            ["Mort. CV", "% Pob. Negra (+)", "Si (UAI movilidad)"],
            ["Hosp. Resp", "% Pob. Negra (+)", "Si (UAI movilidad)"],
        ],
        COLORS['health']
    )

    # SLIDE 10: Cuadrantes
    add_table_slide(
        prs,
        "Los 4 Cuadrantes de Sao Paulo: 645 municipios, 4 realidades",
        ["Cuadrante", "N", "Caracteristica", "Estrategia"],
        [
            ["Q1 - Modelo", "212", "Bajo riesgo, baja vulnerab.", "Mantener"],
            ["Q2 - Conservar", "110", "Alto riesgo, baja vulnerab.", "Fortalecer"],
            ["Q3 - VULNERABLE", "210", "Bajo riesgo, ALTA vulnerab.", "INTERVENCION"],
            ["Q4 - Desarrollo", "113", "Alto riesgo, alta vulnerab.", "Restaurar"],
        ],
        COLORS['vulnerability']
    )

    # SLIDE 11: 10 municipios
    add_table_slide(
        prs,
        "Los 10 Municipios del Workshop: 10 realidades para explorar",
        ["Municipio", "Cuadrante", "Pobreza", "Riesgo", "Dengue", "Gobernanza"],
        [
            ["Iporanga", "Q3", "54.2%", "0.58", "Alta", "Baja"],
            ["Campinas", "Q1", "8.2%", "0.22", "Media", "Alta"],
            ["Santos", "Q1", "9.1%", "0.18", "Media", "Alta"],
            ["S. Joaquim Barra", "Q3", "28.4%", "0.45", "Alta", "Baja"],
            ["Miracatu", "Q3", "38.1%", "0.41", "Alta", "Media"],
            ["Eldorado", "Q4", "43.0%", "0.44", "Alta", "Baja"],
            ["Francisco Morato", "Q4", "35.2%", "0.39", "Alta", "Baja"],
            ["Sao Paulo", "Q1", "11.5%", "0.25", "Alta", "Alta"],
            ["Aruja", "Q2", "15.3%", "0.35", "Media", "Media"],
            ["Cerquilho", "Q2", "12.1%", "0.31", "Baja", "Media"],
        ],
        COLORS['governance']
    )

    # SLIDE 12: Paradojas
    add_content_slide(
        prs,
        "5 Paradojas que desafian la intuicion",
        [
            {'text': "1. Paradoja de gobernanza", 'size': 20, 'bold': True},
            {'text': "   Mas politicas --> mas riesgo (reactivo, no preventivo)", 'size': 16},
            "",
            {'text': "2. Paradoja de pobreza", 'size': 20, 'bold': True},
            {'text': "   Los mas vulnerables --> menos capacidad adaptativa", 'size': 16},
            "",
            {'text': "3. Paradoja del bosque", 'size': 20, 'bold': True},
            {'text': "   Mas cobertura forestal --> menos gobernanza (abandono territorial?)", 'size': 16},
            "",
            {'text': "4. Paradoja de salud", 'size': 20, 'bold': True},
            {'text': "   Mortalidad CV no correlaciona con riesgo climatico (subregistro)", 'size': 16},
            "",
            {'text': "5. Paradoja de polinizacion", 'size': 20, 'bold': True},
            {'text': "   Deficit polinizacion --> mas gobernanza ambiental (proxy de degradacion)", 'size': 16},
        ],
        COLORS['dark']
    )

    # ==========================
    # BLOQUE C: Y AHORA QUE?
    # ==========================

    # SLIDE 13: Actividades
    add_content_slide(
        prs,
        "De la evidencia a la accion: Las actividades de hoy",
        [
            {'text': "ACTIVIDAD 2.1: Ranking de Municipios (30 min)", 'size': 22, 'bold': True, 'color': COLORS['governance']},
            "En grupos, ordenen los 10 municipios por prioridad de intervencion",
            "Justifiquen sus criterios",
            "",
            {'text': "ACTIVIDAD 2.2: Knowledge Budgets (45 min)", 'size': 22, 'bold': True, 'color': COLORS['biodiversity']},
            "Tienen 10 creditos para 'comprar' informacion adicional",
            "Que datos necesitan para decidir mejor?",
            "",
            {'text': "ACTIVIDAD 2.3: Discusion de Paradojas (30 min)", 'size': 22, 'bold': True, 'color': COLORS['climate']},
            "Como disenar politicas proactivas, no reactivas?",
            "Como priorizar Q3 (vulnerables sin desastres)?",
        ],
        COLORS['governance']
    )

    # SLIDE 14: Sistema de créditos
    add_table_slide(
        prs,
        "Sistema de Creditos: Que informacion comprarian?",
        ["Capa de informacion", "Costo", "Lo que revela"],
        [
            ["Mapa de inundaciones", "2 creditos", "Riesgo especifico"],
            ["Cobertura forestal", "2 creditos", "Servicios ecosistemicos"],
            ["Incidencia de dengue", "2 creditos", "Carga de enfermedad"],
            ["% Poblacion negra", "1 credito", "Justicia ambiental"],
            ["% Pobreza", "1 credito", "Vulnerabilidad economica"],
            ["Indice gobernanza (UAI)", "3 creditos", "Capacidad institucional"],
            ["Riqueza de especies", "2 creditos", "Biodiversidad"],
            ["Deficit polinizacion", "2 creditos", "Degradacion agricola"],
            ["Mortalidad CV", "2 creditos", "Impacto en salud"],
            ["Poblacion rural", "1 credito", "Contexto territorial"],
        ],
        COLORS['biodiversity']
    )

    # SLIDE 15: Pregunta final
    add_question_slide(
        prs,
        "Como pasamos de gobernanza\nREACTIVA a PROACTIVA?",
        [
            "1. Invertir ANTES del desastre",
            "   Como identificar municipios Q3? Como justificar inversion sin urgencia?",
            "",
            "2. Priorizar VULNERABILIDAD sobre RIESGO",
            "   Como medir vulnerabilidad? Como comunicar urgencia sin desastre?",
            "",
            "3. Integrar BIODIVERSIDAD en politicas de salud",
            "   Como traducir 'efecto dilucion' a politica? Como valorar servicios ecosistemicos?",
        ]
    )

    # ==========================
    # SLIDES DE BACKUP
    # ==========================

    # BACKUP 1: Metodología
    add_content_slide(
        prs,
        "BACKUP: Metodologia Estadistica",
        [
            {'text': "Metodos utilizados:", 'size': 20, 'bold': True},
            "Correlaciones bivariadas (Pearson)",
            "Modelos lineales mixtos (municipio anidado en microrregion)",
            "Analisis de mediacion (test Sobel, bootstrap)",
            "Analisis de moderacion (interacciones)",
            "SEM - Structural Equation Modeling (semopy)",
            "Seleccion de modelos por AIC (Burnham & Anderson)",
            "",
            {'text': "Validacion:", 'size': 20, 'bold': True},
            "Datos DATASUS validados con TABNET (r > 0.99)",
            "GLM Gamma para distribuciones sesgadas",
            "Bootstrap IC 95% para efectos indirectos",
        ],
        COLORS['dark']
    )

    # BACKUP 2: UAI
    add_table_slide(
        prs,
        "BACKUP: Las 5 dimensiones del Urban Adaptation Index",
        ["Dimension", "Variables", "Que mide"],
        [
            ["Housing", "5 indicadores", "Resiliencia habitacional"],
            ["Environmental", "6 indicadores", "Gestion ambiental"],
            ["Food", "4 indicadores", "Seguridad alimentaria"],
            ["Mobility", "5 indicadores", "Infraestructura transporte"],
            ["Climate Risk", "6 indicadores", "Gestion de riesgos"],
        ],
        COLORS['governance']
    )

    # BACKUP 3: Referencias
    add_content_slide(
        prs,
        "BACKUP: Referencias Cientificas",
        [
            {'text': "Marco conceptual:", 'size': 20, 'bold': True},
            "Levers et al. (2025) - Nexus assessment framework",
            "Barreto et al. (2025) - Planetary health",
            "Keesing et al. (2010) - Dilution effect",
            "Neder et al. (2021) - Urban Adaptation Index",
            "",
            {'text': "Metodologia:", 'size': 20, 'bold': True},
            "Nakagawa & Schielzeth (2013) - R2 marginal",
            "Burnham & Anderson (2002) - Model selection AIC",
            "MacKinnon et al. (2002) - Mediation analysis",
        ],
        COLORS['dark']
    )

    # SLIDE FINAL: Agradecimientos
    add_title_slide(
        prs,
        "Gracias!",
        "Proyecto en colaboracion entre\nUniversidade de Sao Paulo y University of York\n\nContacto: adrian.gonzalez@usp.br"
    )

    # Guardar
    output_path = os.path.join(os.path.dirname(__file__), "Workshop_SEMIL_USP_Day2_MEJORADO.pptx")
    prs.save(output_path)
    print(f"Presentacion guardada en: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()

"""
Script para generar la presentación del Workshop SEMIL-USP
VERSION 2: Con logos y mapas
Autor: Science Team
Fecha: 02/02/2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os

# Rutas
BASE_DIR = os.path.dirname(__file__)
ASSETS_DIR = os.path.join(BASE_DIR, "assets")
FIGURES_DIR = r"C:\Users\arlex\Documents\Adrian David\outputs\figures"

# Colores del proyecto
COLORS = {
    'governance': RGBColor(0x2E, 0x86, 0xAB),
    'biodiversity': RGBColor(0x28, 0xA7, 0x45),
    'climate': RGBColor(0xF5, 0x7C, 0x00),
    'health': RGBColor(0xDC, 0x35, 0x45),
    'vulnerability': RGBColor(0x6F, 0x42, 0xC1),
    'dark': RGBColor(0x2C, 0x3E, 0x50),
    'light': RGBColor(0xEC, 0xF0, 0xF1),
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'black': RGBColor(0x00, 0x00, 0x00),
}

def add_title_slide_with_logos(prs, title, subtitle=""):
    """Agrega slide de título con logos o texto institucional"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fondo
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark']
    bg.line.fill.background()

    # Logo USP (izquierda) - usar texto si no hay imagen
    logo_usp = os.path.join(ASSETS_DIR, "logo_usp.png")
    if os.path.exists(logo_usp):
        try:
            slide.shapes.add_picture(logo_usp, Inches(0.5), Inches(0.3), height=Inches(0.8))
        except:
            pass

    # Texto institucional izquierda
    inst_left = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(3), Inches(0.6))
    tf = inst_left.text_frame
    p = tf.paragraphs[0]
    p.text = "USP"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Logo York (derecha) - usar texto si no hay imagen
    logo_york = os.path.join(ASSETS_DIR, "logo_york.png")
    if os.path.exists(logo_york):
        try:
            slide.shapes.add_picture(logo_york, Inches(12), Inches(0.3), height=Inches(0.8))
        except:
            pass

    # Texto institucional derecha
    inst_right = slide.shapes.add_textbox(Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.6))
    tf = inst_right.text_frame
    p = tf.paragraphs[0]
    p.text = "YORK"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.RIGHT

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.5), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['light']
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(20)

    return slide

def add_content_slide(prs, title, content_items, highlight_color=None):
    """Agrega slide de contenido"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    if highlight_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = highlight_color
        bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.5), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.font.size = Pt(item.get('size', 18))
            p.font.bold = item.get('bold', False)
            if item.get('color'):
                p.font.color.rgb = item['color']
            else:
                p.font.color.rgb = COLORS['dark']
        else:
            p.text = f"  {item}" if item else ""
            p.font.size = Pt(18)
            p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(8)

    return slide

def add_content_slide_with_image(prs, title, content_items, image_path, highlight_color=None):
    """Agrega slide de contenido con imagen a la derecha"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    if highlight_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = highlight_color
        bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    # Contenido a la izquierda
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.font.size = Pt(item.get('size', 16))
            p.font.bold = item.get('bold', False)
            if item.get('color'):
                p.font.color.rgb = item['color']
            else:
                p.font.color.rgb = COLORS['dark']
        else:
            p.text = f"  {item}" if item else ""
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS['dark']
        p.space_before = Pt(6)

    # Imagen a la derecha
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(6.8), Inches(1.4), width=Inches(6))

    return slide

def add_image_slide(prs, title, image_path, caption="", highlight_color=None):
    """Agrega slide con imagen central"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    if highlight_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = highlight_color
        bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    p.alignment = PP_ALIGN.CENTER

    # Imagen centrada
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.5), Inches(1.1), width=Inches(10.5))

    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.5), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_two_images_slide(prs, title, image1_path, image2_path, label1="", label2="", highlight_color=None):
    """Agrega slide con dos imágenes lado a lado"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    if highlight_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = highlight_color
        bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.5), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']
    p.alignment = PP_ALIGN.CENTER

    # Imagen 1 (izquierda)
    if os.path.exists(image1_path):
        slide.shapes.add_picture(image1_path, Inches(0.3), Inches(1.2), width=Inches(6.3))

    # Imagen 2 (derecha)
    if os.path.exists(image2_path):
        slide.shapes.add_picture(image2_path, Inches(6.8), Inches(1.2), width=Inches(6.3))

    # Labels
    if label1:
        lab_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.5), Inches(6.3), Inches(0.4))
        tf = lab_box.text_frame
        p = tf.paragraphs[0]
        p.text = label1
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER

    if label2:
        lab_box = slide.shapes.add_textbox(Inches(6.8), Inches(6.5), Inches(6.3), Inches(0.4))
        tf = lab_box.text_frame
        p = tf.paragraphs[0]
        p.text = label2
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['dark']
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_table_slide(prs, title, headers, rows, highlight_color=None):
    """Agrega slide con tabla"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    if highlight_color:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = highlight_color
        bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark']

    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_width = Inches(12.5)

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.3), table_width, Inches(0.4 * num_rows)).table

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS['dark']
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(12)
        para.font.color.rgb = COLORS['white']
        para.alignment = PP_ALIGN.CENTER

    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = COLORS['dark']
            para.alignment = PP_ALIGN.CENTER
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['light']

    return slide

def add_highlight_slide(prs, main_text, subtext="", number="", color=None):
    """Agrega slide de highlight"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color if color else COLORS['governance']
    bg.line.fill.background()

    if number:
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.5), Inches(1.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(96)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    y_pos = Inches(3) if number else Inches(2.5)
    main_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12.5), Inches(2))
    tf = main_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = main_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if subtext:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.5), Inches(1.5))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtext
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_question_slide(prs, question, options=None):
    """Agrega slide de pregunta"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark']
    bg.line.fill.background()

    q_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.5), Inches(2))
    tf = q_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = question
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    if options:
        opt_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.5), Inches(2.5))
        tf = opt_box.text_frame
        tf.word_wrap = True
        for i, opt in enumerate(options):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = opt
            p.font.size = Pt(22)
            p.font.color.rgb = COLORS['light']
            p.alignment = PP_ALIGN.LEFT
            p.space_before = Pt(12)

    return slide

def create_presentation():
    """Crea la presentación completa con imágenes"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =====================
    # BLOQUE A: EL PROBLEMA
    # =====================

    # SLIDE 1: Título con logos
    add_title_slide_with_logos(
        prs,
        "Riesgos climaticos, salud y perdida de biodiversidad\npredicen politicas de adaptacion...\nexcepto donde mas se necesitan",
        "Evidencia de 645 municipios de Sao Paulo (2010-2019)\n\nWorkshop SEMIL-USP | 24 de febrero 2026\nUniversidade de Sao Paulo & University of York"
    )

    # SLIDE 2: Pregunta central
    add_question_slide(
        prs,
        "La gobernanza climatica en Sao Paulo\nes PROACTIVA o REACTIVA?",
        [
            "PROACTIVA: Anticipa riesgos, invierte antes de crisis",
            "REACTIVA: Responde a desastres, invierte despues de crisis",
            "",
            "Esta es la pregunta que guia nuestro analisis..."
        ]
    )

    # SLIDE 3: Marco Nexus con diagrama
    add_image_slide(
        prs,
        "El Nexo Biodiversidad-Clima-Salud-Gobernanza",
        os.path.join(FIGURES_DIR, "h1_FIG1_causal_panel.png"),
        "Diagrama causal: Biodiversidad -> Servicios Ecosistemicos -> Salud <- Clima <- Gobernanza",
        COLORS['governance']
    )

    # SLIDE 4: Datos
    add_table_slide(
        prs,
        "645 municipios x 104 variables x 10 anos",
        ["Dimension", "Variables", "Fuente"],
        [
            ["Gobernanza", "26 indicadores UAI", "Neder et al. 2021"],
            ["Biodiversidad", "Riqueza, cobertura, polinizacion", "IBGE, MapBiomas"],
            ["Clima", "Inundacion, fuego, estres hidrico", "INPE, ANA"],
            ["Salud", "7 enfermedades + mortalidad", "DATASUS"],
            ["Vulnerabilidad", "Pobreza, ruralidad, demografia", "CadUnico, IBGE"],
        ],
        COLORS['dark']
    )

    # SLIDE 5: Pregunta del workshop
    add_question_slide(
        prs,
        "Si tuvieramos recursos LIMITADOS\npara intervenir...\n\nComo priorizariamos los 645 municipios?",
        [
            "Donde hay mas riesgo climatico?",
            "Donde hay mas vulnerabilidad social?",
            "Donde hay menos gobernanza?",
        ]
    )

    # ==========================
    # BLOQUE B: LO QUE ENCONTRAMOS
    # ==========================

    # SLIDE 6: Hallazgo 1 - Paradoja gobernanza con mapa
    add_content_slide_with_image(
        prs,
        "Hallazgo 1: Mas gobernanza NO significa menos riesgo",
        [
            {'text': "Correlacion POSITIVA (r = +0.29)", 'size': 20, 'bold': True, 'color': COLORS['climate']},
            "",
            "Municipios con MEJOR gobernanza",
            "tienen MAYOR riesgo climatico",
            "",
            {'text': "Por que?", 'size': 18, 'bold': True},
            "La gobernanza responde al riesgo,",
            "no lo previene.",
            "",
            {'text': "Modelo causal:", 'size': 16, 'bold': True},
            "RIESGO -> DESASTRE -> INVERSION",
            "(reactivo, no preventivo)",
        ],
        os.path.join(FIGURES_DIR, "h1_MAP3_bivariate_governance_climate.png"),
        COLORS['climate']
    )

    # SLIDE 7: Hallazgo 2 - Injusticia climática
    add_highlight_slide(
        prs,
        "La pobreza anula el efecto de la gobernanza",
        "27% de la varianza en gobernanza es explicada por pobreza\n\nLos municipios mas pobres tienen menos capacidad adaptativa\nPero son los que mas la necesitan.",
        "27%",
        COLORS['vulnerability']
    )

    # SLIDE 8: Mapa de pobreza y enfermedad
    add_two_images_slide(
        prs,
        "Mapas Bivariados: Evidencia de Injusticia Ambiental",
        os.path.join(FIGURES_DIR, "h1_MAP5_bivariate_poverty_disease.png"),
        os.path.join(FIGURES_DIR, "h1_MAP2_bivariate_governance_biodiv.png"),
        "Pobreza x Carga de Enfermedad",
        "Gobernanza x Biodiversidad",
        COLORS['vulnerability']
    )

    # SLIDE 9: Hallazgo 3 - Efecto dilución con scatter
    add_content_slide_with_image(
        prs,
        "Hallazgo 3: La biodiversidad protege - pero no a todos igual",
        [
            {'text': "Efecto Dilucion confirmado:", 'size': 18, 'bold': True},
            "",
            "Bosque -> Dengue: r = -0.45",
            "Bosque -> Malaria: r = -0.45",
            "Bosque -> Diarrea: r = -0.41",
            "",
            {'text': "PERO el beneficio NO es igual:", 'size': 18, 'bold': True, 'color': COLORS['health']},
            "",
            "Alta pobreza -> Reduce beneficio",
            "Alta % pob. negra -> Reduce beneficio",
            "",
            "Justicia ambiental: el servicio",
            "ecosistemico existe pero no se",
            "distribuye equitativamente",
        ],
        os.path.join(FIGURES_DIR, "h1_scatter_forest_dengue_pobreza.png"),
        COLORS['biodiversity']
    )

    # SLIDE 10: Mapa bosque-dengue
    add_image_slide(
        prs,
        "Mapa Bivariado: Cobertura Forestal y Dengue",
        os.path.join(FIGURES_DIR, "h1_MAP1_bivariate_forest_dengue.png"),
        "Azul: Alto bosque + Bajo dengue | Rojo: Bajo bosque + Alto dengue",
        COLORS['biodiversity']
    )

    # SLIDE 11: Hallazgo 4 - Predictores por enfermedad
    add_table_slide(
        prs,
        "Hallazgo 4: Cada enfermedad tiene su propia firma de predictores",
        ["Enfermedad", "Mejor predictor", "Gobernanza ayuda?"],
        [
            ["Dengue", "% Rural (-)", "Si (indice general)"],
            ["Malaria", "Cobertura forestal (+)", "Si (UAI riesgo)"],
            ["Leptospirosis", "Biodiversidad (-)", "No significativo"],
            ["Leishmaniasis", "Estres hidrico (-)", "No significativo"],
            ["Diarrea", "Riesgo fuego (+)", "No significativo"],
            ["Mort. CV", "% Pob. Negra (+)", "Si (UAI movilidad)"],
            ["Hosp. Resp", "% Pob. Negra (+)", "Si (UAI movilidad)"],
        ],
        COLORS['health']
    )

    # SLIDE 12: Heatmap de selección de modelos
    add_image_slide(
        prs,
        "Seleccion de Modelos: Variables Especificas vs Indices Compuestos",
        os.path.join(FIGURES_DIR, "h1_model_selection_heatmap.png"),
        "75% de las veces, una variable especifica supera al indice compuesto (AIC)",
        COLORS['health']
    )

    # SLIDE 13: Cuadrantes con mapa
    add_content_slide_with_image(
        prs,
        "Los 4 Cuadrantes de Sao Paulo",
        [
            {'text': "Q1 - Modelo (212)", 'size': 16, 'bold': True, 'color': COLORS['biodiversity']},
            "Bajo riesgo, baja vulnerab. -> Mantener",
            "",
            {'text': "Q2 - Conservar (110)", 'size': 16, 'bold': True, 'color': COLORS['governance']},
            "Alto riesgo, baja vulnerab. -> Fortalecer",
            "",
            {'text': "Q3 - VULNERABLE (210)", 'size': 16, 'bold': True, 'color': COLORS['health']},
            "Bajo riesgo, ALTA vulnerab. -> INTERVENCION",
            "",
            {'text': "Q4 - Desarrollo (113)", 'size': 16, 'bold': True, 'color': COLORS['climate']},
            "Alto riesgo, alta vulnerab. -> Restaurar",
            "",
            {'text': "210 municipios en Q3:", 'size': 14, 'bold': True},
            "Invisibles para politicas reactivas",
        ],
        os.path.join(FIGURES_DIR, "h1_MAP2_bivariate_vuln_climate.png"),
        COLORS['vulnerability']
    )

    # SLIDE 14: 10 municipios del workshop
    add_table_slide(
        prs,
        "Los 10 Municipios del Workshop",
        ["Municipio", "Cuadrante", "Pobreza", "Riesgo", "Dengue", "Gobernanza"],
        [
            ["Iporanga", "Q3", "54.2%", "0.58", "Alta", "Baja"],
            ["Campinas", "Q1", "8.2%", "0.22", "Media", "Alta"],
            ["Santos", "Q1", "9.1%", "0.18", "Media", "Alta"],
            ["S. Joaquim Barra", "Q3", "28.4%", "0.45", "Alta", "Baja"],
            ["Miracatu", "Q3", "38.1%", "0.41", "Alta", "Media"],
            ["Eldorado", "Q4", "43.0%", "0.44", "Alta", "Baja"],
            ["Francisco Morato", "Q4", "35.2%", "0.39", "Alta", "Baja"],
            ["Sao Paulo", "Q1", "11.5%", "0.25", "Alta", "Alta"],
            ["Aruja", "Q2", "15.3%", "0.35", "Media", "Media"],
            ["Cerquilho", "Q2", "12.1%", "0.31", "Baja", "Media"],
        ],
        COLORS['governance']
    )

    # SLIDE 15: Paradojas
    add_content_slide(
        prs,
        "5 Paradojas que desafian la intuicion",
        [
            {'text': "1. Paradoja de gobernanza", 'size': 20, 'bold': True},
            {'text': "   Mas politicas -> mas riesgo (reactivo, no preventivo)", 'size': 16},
            "",
            {'text': "2. Paradoja de pobreza", 'size': 20, 'bold': True},
            {'text': "   Los mas vulnerables -> menos capacidad adaptativa", 'size': 16},
            "",
            {'text': "3. Paradoja del bosque", 'size': 20, 'bold': True},
            {'text': "   Mas cobertura forestal -> menos gobernanza", 'size': 16},
            "",
            {'text': "4. Paradoja de salud", 'size': 20, 'bold': True},
            {'text': "   Mortalidad CV no correlaciona con riesgo climatico", 'size': 16},
            "",
            {'text': "5. Paradoja de polinizacion", 'size': 20, 'bold': True},
            {'text': "   Deficit polinizacion -> mas gobernanza ambiental", 'size': 16},
        ],
        COLORS['dark']
    )

    # ==========================
    # BLOQUE C: Y AHORA QUE?
    # ==========================

    # SLIDE 16: Actividades
    add_content_slide(
        prs,
        "De la evidencia a la accion: Las actividades de hoy",
        [
            {'text': "ACTIVIDAD 2.1: Ranking de Municipios (30 min)", 'size': 22, 'bold': True, 'color': COLORS['governance']},
            "En grupos, ordenen los 10 municipios por prioridad",
            "Justifiquen sus criterios",
            "",
            {'text': "ACTIVIDAD 2.2: Knowledge Budgets (45 min)", 'size': 22, 'bold': True, 'color': COLORS['biodiversity']},
            "Tienen 10 creditos para 'comprar' informacion",
            "Que datos necesitan para decidir mejor?",
            "",
            {'text': "ACTIVIDAD 2.3: Discusion de Paradojas (30 min)", 'size': 22, 'bold': True, 'color': COLORS['climate']},
            "Como disenar politicas proactivas?",
            "Como priorizar Q3 (vulnerables sin desastres)?",
        ],
        COLORS['governance']
    )

    # SLIDE 17: Sistema de créditos
    add_table_slide(
        prs,
        "Sistema de Creditos: Que informacion comprarian?",
        ["Capa de informacion", "Costo", "Lo que revela"],
        [
            ["Mapa de inundaciones", "2 creditos", "Riesgo especifico"],
            ["Cobertura forestal", "2 creditos", "Servicios ecosistemicos"],
            ["Incidencia de dengue", "2 creditos", "Carga de enfermedad"],
            ["% Poblacion negra", "1 credito", "Justicia ambiental"],
            ["% Pobreza", "1 credito", "Vulnerabilidad economica"],
            ["Indice gobernanza (UAI)", "3 creditos", "Capacidad institucional"],
            ["Riqueza de especies", "2 creditos", "Biodiversidad"],
            ["Deficit polinizacion", "2 creditos", "Degradacion agricola"],
            ["Mortalidad CV", "2 creditos", "Impacto en salud"],
            ["Poblacion rural", "1 credito", "Contexto territorial"],
        ],
        COLORS['biodiversity']
    )

    # SLIDE 18: Pregunta final
    add_question_slide(
        prs,
        "Como pasamos de gobernanza\nREACTIVA a PROACTIVA?",
        [
            "1. Invertir ANTES del desastre",
            "2. Priorizar VULNERABILIDAD sobre RIESGO",
            "3. Integrar BIODIVERSIDAD en politicas de salud",
        ]
    )

    # ==========================
    # SLIDES DE BACKUP
    # ==========================

    # BACKUP 1: Gobernanza todas las dimensiones
    add_image_slide(
        prs,
        "BACKUP: Gobernanza y todas las dimensiones del Nexus",
        os.path.join(FIGURES_DIR, "h1_governance_all_dimensions.png"),
        "Relaciones entre gobernanza (UAI) y biodiversidad, clima, salud, vulnerabilidad",
        COLORS['governance']
    )

    # BACKUP 2: Componentes UAI
    add_image_slide(
        prs,
        "BACKUP: Componentes del Urban Adaptation Index (UAI)",
        os.path.join(FIGURES_DIR, "h1_governance_components_heatmap.png"),
        "UAI_Movilidad emerge como predictor dominante en casi todas las dimensiones",
        COLORS['governance']
    )

    # BACKUP 3: Heatmap nexus
    add_image_slide(
        prs,
        "BACKUP: Correlaciones del Nexus Completo",
        os.path.join(FIGURES_DIR, "h1_heatmap_nexus.png"),
        "38/51 correlaciones significativas entre biodiversidad, clima, salud y gobernanza",
        COLORS['dark']
    )

    # BACKUP 4: SEM dengue
    add_image_slide(
        prs,
        "BACKUP: Modelo SEM para Dengue",
        os.path.join(FIGURES_DIR, "h1_sem_dengue.png"),
        "Structural Equation Model: Biodiversidad -> Deficit Polinizacion -> Dengue",
        COLORS['health']
    )

    # BACKUP 5: Referencias
    add_content_slide(
        prs,
        "BACKUP: Referencias Cientificas",
        [
            {'text': "Marco conceptual:", 'size': 20, 'bold': True},
            "Levers et al. (2025) - Nexus assessment framework",
            "Barreto et al. (2025) - Planetary health",
            "Keesing et al. (2010) - Dilution effect",
            "Neder et al. (2021) - Urban Adaptation Index",
            "",
            {'text': "Metodologia:", 'size': 20, 'bold': True},
            "Nakagawa & Schielzeth (2013) - R2 marginal",
            "Burnham & Anderson (2002) - Model selection AIC",
            "MacKinnon et al. (2002) - Mediation analysis",
        ],
        COLORS['dark']
    )

    # SLIDE FINAL: Agradecimientos con logos
    add_title_slide_with_logos(
        prs,
        "Obrigado! / Gracias! / Thank you!",
        "Projeto em colaboracao entre\nUniversidade de Sao Paulo & University of York\n\nContato: adrian.gonzalez@usp.br\n\nDados e codigo: github.com/adgch86/saopaulo-biodiversity-health"
    )

    # Guardar
    output_path = os.path.join(BASE_DIR, "Workshop_SEMIL_USP_Day2_v2_con_mapas.pptx")
    prs.save(output_path)
    print(f"Presentacion guardada en: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    return output_path

if __name__ == "__main__":
    create_presentation()
